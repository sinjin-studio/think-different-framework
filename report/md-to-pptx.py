#!/usr/bin/env python3
"""
Convert Think Different Framework markdown presentations to real .pptx files.

Usage: python3 md-to-pptx.py input.md output.pptx
"""

import sys
import re
import os

def check_dependencies():
    try:
        import pptx
        return True
    except ImportError:
        print("python-pptx not installed. Install with: pip3 install python-pptx", file=sys.stderr)
        print("Markdown output is still available.", file=sys.stderr)
        return False


def parse_markdown(text):
    """Parse the presentation markdown into metadata and sections."""
    metadata = {}
    sections = []

    # Parse header metadata (lines like > **Key:** Value)
    for match in re.finditer(r'>\s*\*\*(\w[\w\s]*):\*\*\s*(.+)', text):
        key = match.group(1).strip().lower()
        metadata[key] = match.group(2).strip()

    # Parse title (first # line)
    title_match = re.search(r'^#\s+(.+)', text, re.MULTILINE)
    if title_match:
        metadata['title'] = title_match.group(1).strip()

    # Parse sections: **Section Title** followed by body text
    section_pattern = r'\*\*([^*]+)\*\*\s*\n\n(.*?)(?=\n\*\*[^*]+\*\*|\Z)'
    for match in re.finditer(section_pattern, text, re.DOTALL):
        title = match.group(1).strip()
        body = match.group(2).strip()
        if body:
            sections.append((title, body))

    return metadata, sections


def create_pptx(metadata, sections, output_path):
    """Create a branded .pptx from parsed markdown."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    DARK_BG = RGBColor(0x18, 0x18, 0x18)
    GREEN = RGBColor(0x16, 0x9B, 0x62)
    ORANGE = RGBColor(0xFF, 0x82, 0x00)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GREY = RGBColor(0xCC, 0xCC, 0xCC)

    # Try to find Literata, fall back to Georgia then serif
    HEADING_FONT = 'Literata'
    BODY_FONT = 'Literata'

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def set_slide_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_textbox(slide, left, top, width, height):
        return slide.shapes.add_textbox(left, top, width, height)

    # ── Title slide ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BG)

    # Title
    topic = metadata.get('seed', metadata.get('topic', 'Think Different'))
    tb = add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10.333), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = topic
    p.font.size = Pt(40)
    p.font.color.rgb = WHITE
    p.font.name = HEADING_FONT
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # Subtitle line
    date_str = metadata.get('date', '')
    words_str = metadata.get('words', '')
    subtitle_parts = []
    if date_str:
        subtitle_parts.append(date_str)
    if words_str:
        subtitle_parts.append(f"{words_str} words")

    if subtitle_parts:
        tb2 = add_textbox(slide, Inches(1.5), Inches(4.2), Inches(10.333), Inches(0.6))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = ' / '.join(subtitle_parts)
        p2.font.size = Pt(16)
        p2.font.color.rgb = LIGHT_GREY
        p2.font.name = BODY_FONT
        p2.alignment = PP_ALIGN.LEFT

    # Green accent bar
    from pptx.util import Emu
    bar = slide.shapes.add_shape(
        1, Inches(1.5), Inches(4.0), Inches(2.0), Emu(40000)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()

    # ── Content slides ──
    for title, body in sections:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide, DARK_BG)

        # Section title
        is_sources = 'source' in title.lower()
        title_color = ORANGE if is_sources else GREEN

        tb = add_textbox(slide, Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.color.rgb = title_color
        p.font.name = HEADING_FONT
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

        # Body text - split into paragraphs for readability
        body_top = Inches(1.8)
        body_height = Inches(5.2)
        tb = add_textbox(slide, Inches(1.0), body_top, Inches(11.333), body_height)
        tf = tb.text_frame
        tf.word_wrap = True

        paragraphs = body.split('\n\n')
        for i, para_text in enumerate(paragraphs):
            para_text = para_text.strip()
            if not para_text:
                continue
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
                # Add spacing between paragraphs
                p.space_before = Pt(12)

            p.text = para_text
            p.font.size = Pt(16)
            p.font.color.rgb = WHITE
            p.font.name = BODY_FONT
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = Pt(24)

    prs.save(output_path)


def main():
    if len(sys.argv) != 3:
        print(f"Usage: {sys.argv[0]} input.md output.pptx", file=sys.stderr)
        sys.exit(1)

    input_path = sys.argv[1]
    output_path = sys.argv[2]

    if not check_dependencies():
        sys.exit(1)

    if not os.path.exists(input_path):
        print(f"Error: input file not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    with open(input_path, 'r') as f:
        text = f.read()

    metadata, sections = parse_markdown(text)

    if not sections:
        print("Warning: no sections found in markdown", file=sys.stderr)
        sys.exit(1)

    create_pptx(metadata, sections, output_path)
    print(f"  Created: {output_path}")


if __name__ == '__main__':
    main()
